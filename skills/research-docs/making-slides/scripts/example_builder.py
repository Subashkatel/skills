"""Reference slide builder for the making-slides skill (self-contained).

Demonstrates the full pattern on a real 4-slide deck (a QEC control-loop
simulator example): dual pptx+matplotlib backends, modular type scale,
letterspaced caps titles, ink-only text with color reserved for data marks.
Usage: python example_builder.py [out.pptx] [preview_dir]

Original deck notes follow.

Style follows the user's examples (Dally "Hardware for Deep Learning" Hot Chips
deck, Fan BAICS talk): plain white slides, one centered bold title, ONE focal
element per slide, numbers annotated directly on the diagram, no boxed panels,
kickers, or footers.

Every number on these slides comes from one actual run:

    simulate(RunSpec(ops=cnot_plus_two_t_circuit(),
                     decoder=PerRoundDecoder(tau_us=1.0)), verbose=True)

Trace landmarks (microseconds):
    Op0 CNOT rounds 0.000-6.600 (6 x 1.1), Op1 T 6.600-9.900 (3 x 1.1)
    round hop: fires 1.100 -> controller 1.250 (t_qc) -> decoder 3.250 (t_cd)
    decodes: Op0 W0 8.750-14.750, Op0 W1 15.250-21.250,
             Op1 W0 21.750-27.750 (ready = 21.250 + t_dd 0.5),
             Op2 W0 38.350-44.350
    Op0 result 22.250: Pauli-frame update only (Clifford, no QPU instruction)
    Op1 decision: publish 28.750 (t_do) -> controller 32.750 (t_oc)
                  -> chip 32.900 (t_cq); Op2 starts, stalled since 9.900
    chip_done 36.200, fully_done 45.350
    scoreboard: 4 windows, 20 idle memory rounds, peak 9 payloads buffered

Text is written as native pptx text boxes (crisp at any zoom, editable) and the
same layout is mirrored to high-dpi PNG previews with matplotlib.
"""
from pathlib import Path
import glob
import sys

try:
    from lxml import etree
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.util import Inches, Pt
except ImportError as error:
    raise SystemExit(f"{error}\nInstall the deck dependencies first: "
                     "pip install python-pptx (pulls lxml and Pillow)")

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager
from matplotlib.patches import FancyBboxPatch


def _register_preview_serif():
    """Find a Palatino-class serif for the preview; fall back gracefully.

    The pptx itself just names "Palatino Linotype" — the viewer's machine
    resolves it. Only the matplotlib preview needs a real local font file.
    """
    for pattern in ("/usr/share/fonts/urw-base35/P052-*.otf",     # Linux URW
                    "/usr/share/fonts/**/P052*.otf"):
        files = glob.glob(pattern, recursive=True)
        if files:
            for font_file in files:
                font_manager.fontManager.addfont(font_file)
            return ["P052", "DejaVu Serif"]
    installed = {f.name for f in font_manager.fontManager.ttflist}
    for family in ("Palatino Linotype", "Palatino", "Book Antiqua",
                   "URW Palladio L"):
        if family in installed:
            return [family, "DejaVu Serif"]
    return ["DejaVu Serif"]                       # always present with mpl


matplotlib.rcParams["font.family"] = _register_preview_serif()
matplotlib.rcParams["mathtext.fontset"] = "stix"   # serif math for t-labels

OUT_PPTX = Path(sys.argv[1]) if len(sys.argv) > 1 else Path("slides.pptx")
PREVIEW_DIR = Path(sys.argv[2]) if len(sys.argv) > 2 else Path("slide_previews")

PAGE_W, PAGE_H = 13.333, 7.5

# validated palette (dataviz reference, light mode), slots 1-3 in fixed order
BLUE = "#2a78d6"    # QPU / forward path
AQUA = "#1baf7a"    # decoder
YELLOW = "#eda100"  # feedback / decision
INK = "#0b0b0b"
INK_2 = "#52514e"
MUTED = "#8a887f"
STALL_FILL = "#f1f0ec"
STALL_EDGE = "#9b998f"
WHITE = "#ffffff"

# pastel component fills for the architecture diagram (reference figure)
GREEN_P = "#c9d7b4"
BLUE_P = "#aac5e6"
ORANGE_P = "#e0a26b"
PURPLE_P = "#d4b7cb"

# One versatile family (Bringhurst 6.5.1; Proportional Web §1.1.1): Palatino.
# Office renders Palatino Linotype; Linux substitutes URW P052 by metrics.
FONT = "Palatino Linotype"

# Modular type scale (Bringhurst 3.1.1 "don't compose without a scale";
# Proportional Web fractions of a 12 pt root: 7/8, 1, 9/8, 5/4, 3/2, 2, 5/2, 3).
# Every text size on these slides is one of these steps.
SIZE_XS, SIZE_S, SIZE_M, SIZE_L, SIZE_XL = 10.5, 12, 13.5, 15, 18
SIZE_STAT, SIZE_TITLE = 24, 30
CAPS_TRACKING = 0.07   # letterspacing for strings of capitals (2.1.6: 5-10%)


def _rgb(hex_color):
    return RGBColor.from_string(hex_color.lstrip("#"))


# ----------------------------------------------------------------- the run
#
# The deck is generated from the FULL pipeline: the logical circuit
# (CX q0,q1; T q1; T q1; MZ q1) is compiled by QLX in the container
# (tools/qlx_capture/capture_cnot_then_2t.py -> schedule JSON artifact),
# decsim's qlx_frontend turns the compiled schedule into Operations, and
# build() simulates them and parses every slide number out of the verbose
# trace. Fallbacks: no QLX artifact -> decsim's hand-written quickstart
# ops; no decsim -> a frozen copy of the QLX-compiled run.

import os
import re

DECSIM_DIR = Path(os.environ.get(
    "DECSIM_PATH", Path(__file__).resolve().parents[1] / "decsim"))
QLX_SCHEDULE = Path(os.environ.get(
    "QLX_SCHEDULE",
    Path(__file__).resolve().parent / "qlx_capture"
    / "schedule_cnot_then_2t.json"))

TAU_US = 0.5          # decoder speed: must sustain 6-round windows
NUM_UNITS = 2         # one per distillation stream, shared with compute

FRIENDLY = {"alloc": "alloc", "dealloc": "dealloc", "prep_z": "prep",
            "h": "H", "cx": "CX", "transport": "move T",
            "produce_resource": "distill T", "inject": "T inject",
            "mz": "MZ", "measure_product": "MZZ"}


def display(name):
    """Short display label for a decsim op name (QLX or quickstart style)."""
    if "[" in name:
        return FRIENDLY.get(name.split("[")[0], name.split("[")[0])
    if ":" in name:
        prefix, rest = name.split(":", 1)
        return f"{prefix} {rest.split('(')[0]}"
    return name


def merged_span(intervals):
    """Total length of the union of (start, end) intervals."""
    total, cursor = 0.0, None
    for start, end in sorted(intervals):
        if cursor is None or start > cursor:
            total += end - start
            cursor = end
        elif end > cursor:
            total += end - cursor
            cursor = end
    return total


# snapshot of the last successful live run; the fallback when decsim (or
# the QLX artifact) is unavailable on this machine
FROZEN_PATH = Path(__file__).resolve().parent / "frozen_run.json"


def load_qlx_operations():
    """Compiled schedule artifact -> decsim Operations (+ feedback wiring)."""
    import json

    from decsim.frontends.qlx import qlx_frontend

    payload = json.loads(QLX_SCHEDULE.read_text())
    program = qlx_frontend(payload["entries"])
    injections = [op for op in program.operations if op.consumes_magic_state]
    if len(injections) >= 2:
        # the second T may start only after the first T's decoded outcome
        # (QLX's schedule does not mark classical conditioning; the caller
        # wires it — see qlx_frontend's feedback_candidates contract)
        injections[1].blocked_by = injections[0].id
    blocked = (program.operations[injections[1].blocked_by].name,
               injections[1].name) if len(injections) >= 2 else None
    distill = max((program.rounds.rounds_by_op[op.id]
                   for op in program.operations
                   if op.name.startswith("produce_resource")), default=None)
    cx_rounds = next((program.rounds.rounds_by_op[op.id]
                      for op in program.operations
                      if op.name.startswith("cx")), None)
    return (program.operations, program.rounds, blocked, distill, cx_rounds,
            payload.get("program", "qlx"))


def run_simulation():
    """Compile-from-QLX -> simulate -> extract the slide numbers."""
    import contextlib
    import io

    sys.path.insert(0, str(DECSIM_DIR))
    from decsim.run_spec import RunSpec, simulate
    from decsim.decoders import PerRoundDecoder
    from decsim.config import TimingConfig

    rounds_policy, blocked, distill, cx_rounds = None, None, None, None
    if QLX_SCHEDULE.exists():
        (ops, rounds_policy, blocked, distill, cx_rounds,
         program_name) = load_qlx_operations()
        program_name += " · compiled by QLX"
        source = f"QLX schedule {QLX_SCHEDULE.name} -> decsim at {DECSIM_DIR}"
    else:
        from decsim.frontends.circuit import cnot_plus_two_t_circuit
        ops = cnot_plus_two_t_circuit()
        blocked = (ops[1].name, ops[2].name)
        program_name = "cnot_plus_two_t_circuit (hand-written ops)"
        source = f"decsim quickstart at {DECSIM_DIR} (no QLX artifact)"

    buffer_out = io.StringIO()
    with contextlib.redirect_stdout(buffer_out):
        result = simulate(RunSpec(ops=ops, rounds_policy=rounds_policy,
                                  decoder=PerRoundDecoder(tau_us=TAU_US),
                                  num_units=NUM_UNITS),
                          verbose=True)
    trace = buffer_out.getvalue()
    cluster = result["cluster"]
    cfg = TimingConfig()

    def stamps(pattern):
        return [(float(m.group(1)),) + m.groups()[1:]
                for m in re.finditer(pattern, trace)]

    starts = {op: t for t, op in
              stamps(r"\[\s*([\d.]+) us\] Chip: START (\S+)")}
    dones = {op: t for t, op in
             stamps(r"\[\s*([\d.]+) us\] Chip: (\S+) body done")}
    rounds_of = {op: int(n) for _t, op, n in
                 stamps(r"\[\s*([\d.]+) us\] Chip: (\S+) fires round \d+/(\d+)")}
    op_events = [dict(name=op, label=display(op), start=t,
                      end=dones.get(op, t), rounds=rounds_of.get(op, 1))
                 for op, t in sorted(starts.items(), key=lambda kv: kv[1])]

    decode_starts = stamps(
        r"\[\s*([\d.]+) us\] DecoderCluster: START DECODE (\S+) (W\d+)")
    decode_dones = {(op, w): t for t, op, w in stamps(
        r"\[\s*([\d.]+) us\] DecoderCluster: DECODE DONE (\S+) (W\d+)")}
    windows = [dict(name=op, w=w, label=f"{display(op)} · {w}",
                    start=t, end=decode_dones[(op, w)])
               for t, op, w in decode_starts]

    def take(pattern):
        found = stamps(pattern)
        return found[0][0] if found else None

    frame = stamps(r"\[\s*([\d.]+) us\] Orchestrator: result for (\S+): "
                   r"Pauli-frame update")
    return dict(
        source=source, program=program_name,
        distance=cluster.code.distance, round_us=cfg.round_us,
        commit=cluster.commit, buffer=cluster.buffer, tau_us=TAU_US,
        num_units=NUM_UNITS,
        timing=dict(qc=cfg.t_qc_us, cd=cfg.t_cd_us, dd=cfg.t_dd_us,
                    do=cfg.t_do_us, oc=cfg.t_oc_us, cq=cfg.t_cq_us),
        ops=op_events, windows=windows, blocked=blocked,
        distill_rounds=distill, cx_rounds=cx_rounds,
        first_fire=take(r"\[\s*([\d.]+) us\] Chip: \S+ fires round 1/"),
        first_arrival=take(
            r"\[\s*([\d.]+) us\] DecoderCluster: round 1 of"),
        first_frame=frame[0] if frame else None,
        decision=take(
            r"\[\s*([\d.]+) us\] Orchestrator: DISPATCH conditional release"),
        ctrl_release=take(
            r"\[\s*([\d.]+) us\] Controller: received release instruction"),
        consumed=take(r"\[\s*([\d.]+) us\] Chip: CONSUMED decision"),
        chip_done=result["chip_done"] / 1e6,
        fully_done=result["fully_done"] / 1e6,
        memory_rounds=cluster.memory_rounds_total,
        peak_payloads=cluster.peak_payloads,
        total_windows=cluster.total_windows,
    )


def get_data():
    import json
    try:
        data = run_simulation()
    except Exception as error:
        if not FROZEN_PATH.exists():
            raise SystemExit(
                f"cannot run the simulation ({error}) and no frozen "
                f"snapshot exists at {FROZEN_PATH}")
        print(f"note: falling back to the frozen snapshot ({error})")
        data = json.loads(FROZEN_PATH.read_text())
        data["source"] = f"frozen snapshot {FROZEN_PATH.name}"
        return data
    FROZEN_PATH.write_text(json.dumps(data, indent=1))
    return data


def dur(x):
    """Duration formatting: 6.0, 0.15, 11.85, 23.0."""
    s = f"{x:.2f}".rstrip("0").rstrip(".")
    return s if "." in s else s + ".0"


def us3(x):
    """Timestamp formatting: 32.900."""
    return f"{x:.3f}"


class SlideCanvas:
    """One slide drawn to both backends: pptx shapes and a matplotlib mirror."""

    def __init__(self, pptx_slide, mpl_ax):
        self.shapes = pptx_slide.shapes
        self.ax = mpl_ax
        self._z = 2.0

    def _next_z(self):
        """Monotonic zorder so matplotlib paints in call order, like pptx."""
        self._z += 0.01
        return self._z

    # ------------------------------------------------------------ primitives

    def box(self, x, y, w, h, fill=WHITE, edge=None, edge_w=1.25,
            radius=0.05, dash=None):
        shape = self.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
            Inches(w), Inches(h))
        shape.adjustments[0] = min(0.5, radius / min(w, h))
        shape.shadow.inherit = False
        if fill is None:
            shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb(fill)
        if edge is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = _rgb(edge)
            shape.line.width = Pt(edge_w)
            if dash:
                ln = shape.line._get_or_add_ln()
                d = etree.SubElement(ln, qn("a:prstDash"))
                d.set("val", "dash")
        shape.text_frame.paragraphs[0].text = ""

        patch = FancyBboxPatch(
            (x, y), w, h,
            boxstyle=f"round,pad=0,rounding_size={radius}",
            facecolor=fill if fill else "none",
            edgecolor=edge if edge else "none",
            linewidth=edge_w, linestyle="--" if dash else "-",
            zorder=self._next_z())
        self.ax.add_patch(patch)

    def text(self, x, y, w, h, content, size, color=INK, bold=False,
             italic=False, align="left", valign="top", leading=1.2,
             tracking=None, rot=0):
        tb = self.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        # no auto-wrap: line breaks are explicit (\n), so layout cannot drift
        # when the viewer's font metrics differ from the preview face
        tf.word_wrap = False
        if rot:
            tb.rotation = -rot          # pptx rotates clockwise-positive
        for side in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(tf, side, 0)
        tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                              "bottom": MSO_ANCHOR.BOTTOM}[valign]
        pptx_align = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                      "right": PP_ALIGN.RIGHT}[align]
        for i, line in enumerate(content.split("\n")):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.alignment = pptx_align
            para.line_spacing = leading
            run = para.add_run()
            run.text = line
            font = run.font
            font.name = FONT
            font.size = Pt(size)
            font.bold = bold
            font.italic = italic
            font.color.rgb = _rgb(color)
            if tracking:
                # letterspacing, in hundredths of a point
                rPr = run._r.get_or_add_rPr()
                rPr.set("spc", str(int(round(tracking * size * 100))))

        mpl_x = {"left": x, "center": x + w / 2, "right": x + w}[align]
        mpl_y = {"top": y, "middle": y + h / 2, "bottom": y + h}[valign]
        mpl_content = content
        if tracking:
            # approximate letterspacing with hair spaces in the preview
            mpl_content = "\n".join(" ".join(line)
                                    for line in content.split("\n"))
        self.ax.text(mpl_x, mpl_y, mpl_content, fontsize=size, color=color,
                     weight="bold" if bold else "normal",
                     style="italic" if italic else "normal",
                     ha=align, va={"top": "top", "middle": "center",
                                   "bottom": "bottom"}[valign],
                     linespacing=leading + 0.15,
                     rotation=rot, rotation_mode="anchor")

    def line(self, x1, y1, x2, y2, color, width=1.5, dash=None, arrow=False):
        conn = self.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
            Inches(x2), Inches(y2))
        conn.shadow.inherit = False
        conn.line.color.rgb = _rgb(color)
        conn.line.width = Pt(width)
        ln = conn.line._get_or_add_ln()
        if dash:
            d = etree.SubElement(ln, qn("a:prstDash"))
            d.set("val", "dash")
        if arrow:
            tail = etree.SubElement(ln, qn("a:tailEnd"))
            tail.set("type", "triangle")
            tail.set("w", "med")
            tail.set("len", "med")

        z = self._next_z()
        if arrow:
            annotation = self.ax.annotate(
                "", xy=(x2, y2), xytext=(x1, y1),
                arrowprops=dict(arrowstyle="-|>", color=color, lw=width,
                                shrinkA=0, shrinkB=0, mutation_scale=14,
                                linestyle="--" if dash else "-"))
            annotation.arrow_patch.set_zorder(z)
        else:
            self.ax.plot([x1, x2], [y1, y2], color=color, lw=width,
                         linestyle="--" if dash else "-",
                         solid_capstyle="butt", zorder=z)

    def arrow(self, x1, y1, x2, y2, color, width=1.75):
        self.line(x1, y1, x2, y2, color, width, arrow=True)

    def double_arrow(self, x1, y1, x2, y2, color=INK, width=1.25):
        conn = self.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
            Inches(x2), Inches(y2))
        conn.shadow.inherit = False
        conn.line.color.rgb = _rgb(color)
        conn.line.width = Pt(width)
        ln = conn.line._get_or_add_ln()
        for end in ("a:headEnd", "a:tailEnd"):
            e = etree.SubElement(ln, qn(end))
            e.set("type", "triangle")
            e.set("w", "med")
            e.set("len", "med")
        annotation = self.ax.annotate(
            "", xy=(x2, y2), xytext=(x1, y1),
            arrowprops=dict(arrowstyle="<|-|>", color=color, lw=width,
                            shrinkA=0, shrinkB=0, mutation_scale=11))
        annotation.arrow_patch.set_zorder(self._next_z())

    def circle(self, cx, cy, r, fill=None, edge=INK, edge_w=1.5):
        shape = self.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(cx - r), Inches(cy - r),
            Inches(2 * r), Inches(2 * r))
        shape.shadow.inherit = False
        if fill is None:
            shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb(fill)
        if edge is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = _rgb(edge)
            shape.line.width = Pt(edge_w)
        patch = plt.Circle((cx, cy), r,
                           facecolor=fill if fill else "none",
                           edgecolor=edge if edge else "none",
                           linewidth=edge_w, zorder=self._next_z())
        self.ax.add_patch(patch)

    def tlabel(self, cx, cy, sub, rest="", size=SIZE_S, rot=0, color=INK,
               base="t"):
        """Math-style label 'base_sub: rest' centered on (cx, cy)."""
        w, h = 4.2, 0.3
        tb = self.shapes.add_textbox(Inches(cx - w / 2), Inches(cy - h / 2),
                                     Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = False
        for side in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(tf, side, 0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        if rot:
            tb.rotation = -rot
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        para.line_spacing = 1.0
        pieces = [(base, size, True, None),
                  (sub, size * 0.7, False, "-25000")]
        if rest:
            pieces.append((f":  {rest}", size, False, None))
        for content, run_size, run_italic, baseline in pieces:
            run = para.add_run()
            run.text = content
            run.font.name = FONT
            run.font.size = Pt(run_size)
            run.font.italic = run_italic
            run.font.color.rgb = _rgb(color)
            if baseline:
                run._r.get_or_add_rPr().set("baseline", baseline)

        mathtext = f"${base}_{{\\mathrm{{{sub}}}}}$"
        if rest:
            mathtext += f":  {rest}"
        self.ax.text(cx, cy, mathtext, fontsize=size, color=color,
                     ha="center", va="center",
                     rotation=rot, rotation_mode="anchor")

    def elbow_arrow(self, points, color, width=1.75):
        """Polyline with an arrowhead on the final segment."""
        for (ax_, ay), (bx, by) in zip(points[:-2], points[1:-1]):
            self.line(ax_, ay, bx, by, color, width)
        (px, py), (qx, qy) = points[-2], points[-1]
        self.arrow(px, py, qx, qy, color, width)

    def title(self, text_value):
        """Letterspaced capitals in the text weight (2.1.6, 3.5.1: weight
        decreases as size increases). Pass the string already in capitals —
        upper() would corrupt µ."""
        self.text(0.7, 0.5, 11.93, 0.55, text_value, SIZE_STAT, INK,
                  align="center", tracking=CAPS_TRACKING)

    def component(self, x, y, w, h, name, fill, sub=None, text_color=WHITE,
                  name_size=SIZE_L):
        """Dally-style flat colored box with a centered label."""
        self.box(x, y, w, h, fill=fill, edge=None, radius=0.06)
        if sub:
            self.text(x, y + 0.15, w, 0.4, name, name_size, text_color,
                      bold=True, align="center")
            self.text(x, y + h - 0.42, w, 0.3, sub, SIZE_XS, text_color,
                      align="center")
        else:
            self.text(x, y, w, h, name, name_size, text_color, bold=True,
                      align="center", valign="middle")


# ---------------------------------------------------------------- slide 1

def slide_example(c, d):
    c.title("THE EXAMPLE: A CNOT AND TWO DEPENDENT T GATES")

    # the circuit itself: two qubit wires, a CNOT, and two T gates on q1
    WIRE_X0, WIRE_X1 = 2.4, 11.0
    Q0_Y, Q1_Y = 2.55, 3.45
    CNOT_X, T1_X, T2_X = 4.2, 6.7, 9.2

    c.line(WIRE_X0, Q0_Y, WIRE_X1, Q0_Y, INK, 1.5)
    c.line(WIRE_X0, Q1_Y, WIRE_X1, Q1_Y, INK, 1.5)
    c.tlabel(2.0, Q0_Y, "0", base="q", size=SIZE_L)
    c.tlabel(2.0, Q1_Y, "1", base="q", size=SIZE_L)

    # CNOT: control on q0, target on q1
    c.circle(CNOT_X, Q0_Y, 0.06, fill=INK, edge=None)
    c.circle(CNOT_X, Q1_Y, 0.15, fill=None, edge=INK, edge_w=1.5)
    c.line(CNOT_X, Q0_Y, CNOT_X, Q1_Y + 0.15, INK, 1.5)
    c.line(CNOT_X - 0.15, Q1_Y, CNOT_X + 0.15, Q1_Y, INK, 1.5)

    # T gates on q1
    for gate_x in (T1_X, T2_X):
        c.box(gate_x - 0.25, Q1_Y - 0.25, 0.5, 0.5, fill=WHITE, edge=INK,
              edge_w=1.5, radius=0.02)
        c.text(gate_x - 0.25, Q1_Y - 0.25, 0.5, 0.5, "T", SIZE_L, INK,
               italic=True, align="center", valign="middle")

    # classical dependency: Op2 waits for Op1's decoded outcome
    c.line(T1_X, Q1_Y + 0.25, T1_X, 4.15, INK, 1.25, dash=True)
    c.line(T1_X, 4.15, T2_X, 4.15, INK, 1.25, dash=True)
    c.line(T2_X, 4.15, T2_X, Q1_Y + 0.27, INK, 1.25, dash=True, arrow=True)
    c.text(T1_X - 1.0, 4.28, T2_X - T1_X + 2.0, 0.28,
           "decoded outcome returns through the control loop", SIZE_S, INK_2,
           italic=True, align="center")

    if d.get("distill_rounds"):
        t_note = (f"consumes a T state\nfrom a {d['distill_rounds']}-round\n"
                  "15-to-1 distillation")
        cx_note = (f"lattice-surgery CX\n{d['cx_rounds']} round in the\n"
                   "QLX schedule")
    else:
        t_note = "non-Clifford\nneeds a magic state"
        cx_note = "Clifford\nsyndrome rounds"
    notes = [(CNOT_X, cx_note), (T1_X, t_note),
             (T2_X, "may start only after\nthe first outcome\nis decoded")]
    for gate_x, note in notes:
        c.text(gate_x - 1.2, 4.85, 2.4, 0.85, note, SIZE_M, INK_2,
               align="center")

    window_us = (d["commit"] + d["buffer"]) * d["tau_us"]
    c.text(0.7, 6.0, 11.93, 0.7,
           f"d = {d['distance']} surface code   ·   QEC round "
           f"{dur(d['round_us'])} µs   ·   sliding window: commit "
           f"{d['commit']} + buffer {d['buffer']} rounds\n"
           f"{d['num_units']} decoder units   ·   τ = {dur(d['tau_us'])} µs "
           f"per round  →  {dur(window_us)} µs per window",
           SIZE_L, INK, align="center")

    c.text(0.7, 7.0, 11.93, 0.3,
           "qlx.estimate.schedule(cnot_then_2t)  →  "
           "schedule_cnot_then_2t.json  →  decsim.frontends.qlx  →  "
           "simulate(RunSpec(...))",
           SIZE_XS, MUTED, italic=True, align="center")


# ---------------------------------------------------------------- slide 2

def slide_diagram(c, d):
    c.title("THE SIMULATOR: ONE TIMED CONTROL LOOP")
    t = d["timing"]

    # pastel components (reference architecture figure)
    orch = (0.9, 2.3, 2.15, 0.95)
    ctrl = (6.9, 2.3, 1.8, 0.95)
    qpu = (10.9, 2.15, 1.73, 1.25)
    clus = (4.55, 4.2, 2.5, 2.4)
    c.box(*orch, fill=GREEN_P, edge=None, radius=0.12)
    c.text(orch[0], orch[1], orch[2], orch[3], "Orchestrator", SIZE_L, INK,
           align="center", valign="middle")
    c.box(*ctrl, fill=BLUE_P, edge=None, radius=0.12)
    c.text(ctrl[0], ctrl[1], ctrl[2], ctrl[3], "Controller", SIZE_L, INK,
           align="center", valign="middle")
    c.box(*qpu, fill=ORANGE_P, edge=None, radius=0.12)
    c.text(qpu[0], qpu[1], qpu[2], qpu[3], "Quantum\nprocessing unit",
           SIZE_L, INK, align="center", valign="middle")
    c.box(*clus, fill=PURPLE_P, edge=None, radius=0.12)
    c.text(clus[0], clus[1] + 0.12, clus[2], 0.35, "Decoder cluster",
           SIZE_L, INK, align="center")

    # workload feeds the orchestrator
    c.text(0.9, 1.25, 2.15, 0.28, "operation stream", SIZE_S, INK,
           align="center")
    c.arrow(1.97, 1.55, 1.97, 2.28, INK, 1.5)

    # orchestrator -> controller
    c.arrow(3.05, 2.62, 6.88, 2.62, INK, 1.5)
    c.tlabel(4.96, 2.4, "oc", f"decisions · {dur(t['oc'])} µs")

    # controller <-> QPU
    c.arrow(8.72, 2.5, 10.88, 2.5, INK, 1.5)
    c.tlabel(9.8, 2.25, "cq", f"control signals · {dur(t['cq'])} µs",
             size=SIZE_XS)
    c.arrow(10.88, 2.95, 8.74, 2.95, INK, 1.5)
    c.tlabel(9.8, 3.22, "qc", f"syndromes · {dur(t['qc'])} µs", size=SIZE_XS)

    # controller -> decoder cluster (syndromes travel down-left)
    c.arrow(7.7, 3.27, 6.82, 4.18, INK, 1.5)
    c.tlabel(7.85, 4.02, "cd", f"syndromes · {dur(t['cd'])} µs", rot=46)

    # orchestrator <-> decoder cluster (jobs down, results up)
    c.arrow(2.72, 3.27, 4.68, 4.52, INK, 1.5)
    c.text(2.9, 3.45, 2.2, 0.28, "decoding jobs", SIZE_S, INK,
           align="center", valign="middle", rot=-31)
    c.arrow(4.5, 4.75, 2.54, 3.5, INK, 1.5)
    c.tlabel(3.15, 4.45, "do", f"results · {dur(t['do'])} µs", rot=-31)

    # decode units inside the cluster, exchanging window boundaries
    sq = 0.42
    units = [(5.0, 4.85), (6.18, 4.85), (5.0, 5.95), (6.18, 5.95)]
    for ux, uy in units:
        c.box(ux, uy, sq, sq, fill=WHITE, edge=INK, edge_w=1.0, radius=0.02)
    c.double_arrow(5.46, 5.06, 6.14, 5.06)
    c.double_arrow(5.46, 6.16, 6.14, 6.16)
    c.double_arrow(5.21, 5.31, 5.21, 5.91)
    c.double_arrow(6.39, 5.31, 6.39, 5.91)
    c.tlabel(5.8, 5.61, "dd", size=SIZE_S)

    c.tlabel(9.55, 5.0, "dd", f"window → window handoff · {dur(t['dd'])} µs")
    c.text(7.35, 5.28, 4.4, 0.28,
           "(dependent windows wait for their predecessor’s result)",
           SIZE_XS, INK_2, align="center")

    # the example's actual timestamps, directly under the loop
    blocker = next(o for o in d["ops"] if o["name"] == d["blocked"][0])
    final = max((w for w in d["windows"] if w["name"] == blocker["name"]),
                key=lambda w: w["start"])
    c.text(0.7, 6.8, 11.93, 0.6,
           f"In the run:  a round fired at {us3(d['first_fire'])} reaches a "
           f"decode unit at {us3(d['first_arrival'])}."
           f"\nThe first T’s outcome:  decoded {us3(final['end'])}  →  "
           f"orchestrator {us3(d['decision'])}  →  controller "
           f"{us3(d['ctrl_release'])}  →  QPU {us3(d['consumed'])}.",
           SIZE_M, INK_2, align="center")


# ---------------------------------------------------------------- slide 3

def slide_timeline(c, d):
    c.title(f"THE RUN: 0 → {dur(d['fully_done'])} µs")

    X0 = 1.15                      # timeline origin (inches)
    axis_max = (int(d["fully_done"] / 5) + 1) * 5
    us_per_in = (axis_max + 1) / 11.6
    tx = lambda t_us: X0 + t_us / us_per_in

    LANES = [("QPU", 2.15), ("DECODER", 3.65), ("FEEDBACK", 5.05)]
    BAR_H = 0.42

    step = next(s for s in (1, 2, 5, 10, 20, 25, 50)
                if axis_max / s <= 12)
    axis_y = 6.0
    last_tick = 0
    for t in range(0, axis_max + 1, step):
        x = tx(t)
        c.line(x, 1.85, x, axis_y, "#e9e7e0", 0.75)
        c.text(x - 0.3, axis_y + 0.06, 0.6, 0.25, str(t), SIZE_XS, MUTED,
               align="center")
        last_tick = t
    c.text(tx(last_tick) + 0.35, axis_y + 0.06, 0.7, 0.25, "µs", SIZE_XS,
           MUTED)
    for label, y in LANES:
        c.text(0.02, y, 1.02, BAR_H, label, SIZE_S, INK,
               align="right", valign="middle", tracking=CAPS_TRACKING)

    # --- QPU lane: stall first (so op bars paint over its edges)
    y = LANES[0][1]
    stall = None
    if d["blocked"]:
        blocker = next(o for o in d["ops"] if o["name"] == d["blocked"][0])
        waiter = next(o for o in d["ops"] if o["name"] == d["blocked"][1])
        stall = waiter["start"] - blocker["end"]
        c.box(tx(blocker["end"]) + 0.015, y,
              tx(waiter["start"]) - tx(blocker["end"]) - 0.03,
              BAR_H, fill=STALL_FILL, edge=STALL_EDGE, edge_w=1.0, dash=True,
              radius=0.04)
        c.text(tx(blocker["end"]), y + 0.08,
               tx(waiter["start"]) - tx(blocker["end"]), 0.28,
               f"stalled {dur(stall)} µs", SIZE_S, INK_2, align="center",
               italic=True)
    label_slots = []
    for op in d["ops"]:
        w = tx(op["end"]) - tx(op["start"])
        c.box(tx(op["start"]) + 0.015, y, max(w - 0.03, 0.035), BAR_H,
              fill=BLUE, edge=None, radius=0.04)
        mid = tx((op["start"] + op["end"]) / 2)
        # label wide bars inside; label narrow bars below, skipping collisions
        if w > 0.9:
            c.text(mid - w / 2, y + 0.08, w, 0.28, op["label"], SIZE_S,
                   WHITE, bold=True, align="center")
        elif all(abs(mid - taken) > 0.62 for taken in label_slots):
            label_slots.append(mid)
            c.text(mid - 0.6, y + BAR_H + 0.06, 1.2, 0.25, op["label"],
                   SIZE_XS, INK_2, align="center")

    # --- decoder lane
    y = LANES[1][1]
    for win in d["windows"]:
        w = tx(win["end"]) - tx(win["start"])
        c.box(tx(win["start"]) + 0.015, y, max(w - 0.03, 0.035), BAR_H,
              fill=AQUA, edge=None, radius=0.04)
        if w > 1.15:
            c.text(tx(win["start"]), y + 0.08, w, 0.28, win["label"], SIZE_S,
                   WHITE, bold=True, align="center")
    window_rounds = d["commit"] + d["buffer"]
    window_us = window_rounds * d["tau_us"]
    c.text(tx(0), y - 0.30, 5.5, 0.25,
           f"{len(d['windows'])} windows · each {window_rounds} rounds × τ "
           f"= {dur(window_us)} µs · {d['num_units']} units",
           SIZE_S, INK_2)

    # --- feedback lane
    y = LANES[2][1]
    if d["decision"] and d["consumed"]:
        c.box(tx(d["decision"]) + 0.015, y,
              tx(d["consumed"]) - tx(d["decision"]) - 0.03, BAR_H,
              fill=YELLOW, edge=None, radius=0.04)
        c.text(tx(d["consumed"]) - 8.0, y - 0.30, 7.9, 0.25,
               f"decision travels {us3(d['decision'])} → "
               f"{us3(d['consumed'])}",
               SIZE_S, INK, bold=True, align="right")
    if d["first_frame"]:
        frame_t = d["first_frame"][0]
        c.line(tx(frame_t), y - 0.04, tx(frame_t), y + BAR_H + 0.04,
               YELLOW, 2.5)
        note = (f"{us3(frame_t)} · first Clifford result:\n"
                "Pauli-frame update only")
        if tx(frame_t) - 3.3 > 0.1:
            c.text(tx(frame_t) - 3.3, y + 0.52, 3.1, 0.5, note, SIZE_S,
                   INK_2, align="right")
        else:
            c.text(tx(frame_t) + 0.15, y + 0.52, 3.1, 0.5, note, SIZE_S,
                   INK_2)

    # markers (the colored dashed lines carry identity; the labels stay in ink)
    c.line(tx(d["chip_done"]), 1.85, tx(d["chip_done"]), axis_y, BLUE, 1.5,
           dash=True)
    c.text(tx(d["chip_done"]) - 2.35, 1.55, 2.2, 0.25,
           f"QPU done {us3(d['chip_done'])}", SIZE_S, INK,
           bold=True, align="right")
    c.line(tx(d["fully_done"]), 1.85, tx(d["fully_done"]), axis_y, AQUA, 1.5,
           dash=True)
    c.text(tx(d["fully_done"]) - 1.95, 1.25, 1.9, 0.25,
           f"all decoded {us3(d['fully_done'])}", SIZE_S,
           INK, bold=True, align="right")

    busy = merged_span([(o["start"], o["end"]) for o in d["ops"]])
    summary = (f"The QPU computes for {dur(busy)} of its "
               f"{dur(d['chip_done'])} µs")
    if stall is not None:
        summary += (f" — the T-gate feedback stall is {dur(stall)} µs of "
                    "that schedule.")
    else:
        summary += "."
    c.text(0.7, 6.75, 11.93, 0.35, summary, SIZE_L, INK, italic=True,
           align="center")


# ---------------------------------------------------------------- slide 4

def slide_stall(c, d):
    blocker = next(o for o in d["ops"] if o["name"] == d["blocked"][0])
    waiter = next(o for o in d["ops"] if o["name"] == d["blocked"][1])
    # the decision-triggering decode: the blocker's final window
    final = max((w for w in d["windows"] if w["name"] == blocker["name"]),
                key=lambda w: w["start"])
    op1_done, op2_start = blocker["end"], waiter["start"]
    stall = op2_start - op1_done
    wait = final["start"] - op1_done
    decode = final["end"] - final["start"]
    t_do = d["decision"] - final["end"]
    t_oc = d["ctrl_release"] - d["decision"]
    t_cq = d["consumed"] - d["ctrl_release"]

    c.title(f"WHERE THE {dur(stall)} µs STALL GOES")

    scale = 11.0 / stall
    bx0, by, bh = 1.15, 2.2, 0.55
    segments = [(wait, STALL_FILL, STALL_EDGE), (decode, AQUA, None),
                (t_do, YELLOW, None), (t_oc, YELLOW, None),
                (t_cq, YELLOW, None)]
    x = bx0
    seg_x = []
    for width_us, fill, edge in segments:
        w = width_us * scale
        c.box(x + 0.012, by, w - 0.024, bh, fill=fill, edge=edge,
              edge_w=1.0, radius=0.04)
        seg_x.append((x, w))
        x += w

    c.text(seg_x[0][0], by + 0.15, seg_x[0][1], 0.3,
           f"waiting  {dur(wait)} µs", SIZE_M, INK_2, align="center")
    decode_mid = seg_x[1][0] + seg_x[1][1] / 2
    if seg_x[1][1] > 1.4:
        c.text(seg_x[1][0], by + 0.15, seg_x[1][1], 0.3,
               f"decode  {dur(decode)} µs", SIZE_M, WHITE, bold=True,
               align="center")
    else:
        c.line(decode_mid, by - 0.28, decode_mid, by - 0.04, MUTED, 1.0)
        c.text(decode_mid - 1.1, by - 0.56, 2.2, 0.25,
               f"decode  {dur(decode)}", SIZE_S, INK_2, align="center")
    c.text(seg_x[3][0], by + 0.15, seg_x[3][1], 0.3, f"t_oc  {dur(t_oc)}",
           SIZE_M, INK, bold=True, align="center")
    mid2 = seg_x[2][0] + seg_x[2][1] / 2
    c.line(mid2, by - 0.28, mid2, by - 0.04, MUTED, 1.0)
    c.text(mid2 - 1.1, by - 0.56, 2.2, 0.25, f"t_do  {dur(t_do)}", SIZE_S,
           INK_2, align="center")
    mid4 = seg_x[4][0] + seg_x[4][1] / 2
    c.line(mid4, by - 0.28, mid4, by - 0.04, MUTED, 1.0)
    c.text(mid4 - 1.1, by - 0.56, 2.2, 0.25, f"t_cq  {dur(t_cq)}", SIZE_S,
           INK_2, align="center")
    c.text(bx0, by + bh + 0.12, 3.4, 0.25,
           f"{us3(op1_done)} µs · first T done", SIZE_S, INK_2)
    c.text(bx0 + 11.0 - 3.4, by + bh + 0.12, 3.4, 0.25,
           f"{us3(op2_start)} µs · second T starts", SIZE_S, INK_2,
           align="right")

    # sparse breakdown with the numbers (Dally-style)
    window_rounds = d["commit"] + d["buffer"]
    comm = t_do + t_oc + t_cq
    bullets = [
        (f"{dur(wait)} µs",
         "waiting — the first T’s stream pads its decode window "
         f"(memory rounds) and queues for a unit"),
        (f"{dur(decode)} µs",
         f"decoding the first T’s final window ({window_rounds} rounds × τ "
         f"{dur(d['tau_us'])} µs)"),
        (f"{dur(comm)} µs",
         f"communication — publish {dur(t_do)} (t_do) + decision return "
         f"{dur(t_oc)} (t_oc) + delivery {dur(t_cq)} (t_cq)"),
    ]
    y = 4.0
    for number, caption in bullets:
        c.text(1.7, y, 1.7, 0.35, number, SIZE_XL, INK, bold=True,
               align="right")
        c.text(3.65, y + 0.07, 8.5, 0.35, caption, SIZE_M, INK_2)
        y += 0.6

    c.text(0.7, 6.0, 11.93, 0.35,
           f"Whole run:  {d['total_windows']} windows decoded  ·  "
           f"{d['memory_rounds']} idle memory rounds emitted  ·  peak "
           f"{d['peak_payloads']} syndrome payloads buffered",
           SIZE_M, INK_2, align="center")
    c.text(0.7, 6.6, 11.93, 0.4,
           "The feedback loop, not decoder throughput, sets the pace.",
           SIZE_XL, INK, italic=True, align="center")


# ------------------------------------------------------------------- build

def build():
    data = get_data()
    print(f"slide numbers from: {data['source']}")
    PREVIEW_DIR.mkdir(exist_ok=True)
    presentation = Presentation()
    presentation.slide_width = Inches(PAGE_W)
    presentation.slide_height = Inches(PAGE_H)
    blank_layout = presentation.slide_layouts[6]

    slide_fns = [slide_example, slide_diagram, slide_timeline, slide_stall]
    for index, draw in enumerate(slide_fns, 1):
        pptx_slide = presentation.slides.add_slide(blank_layout)
        fig = plt.figure(figsize=(PAGE_W, PAGE_H))
        ax = fig.add_axes([0, 0, 1, 1])
        ax.set_xlim(0, PAGE_W)
        ax.set_ylim(PAGE_H, 0)
        ax.axis("off")
        draw(SlideCanvas(pptx_slide, ax), data)
        fig.savefig(PREVIEW_DIR / f"slide_{index}.png", dpi=180,
                    facecolor="white")
        plt.close(fig)

    presentation.core_properties.title = "decsim — simulator example"
    presentation.save(OUT_PPTX)
    print(OUT_PPTX)
    print(PREVIEW_DIR)


if __name__ == "__main__":
    build()
